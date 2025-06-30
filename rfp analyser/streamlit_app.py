import streamlit as st
import requests
import tempfile
from pathlib import Path
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import base64

# Constants
API_URL = "http://localhost:8000"

def create_pdf(summary_data):
    """Create a PDF from the summary data."""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30
    )
    story.append(Paragraph("RFP Analysis Summary", title_style))
    story.append(Spacer(1, 20))

    # Content
    for section, content in summary_data.items():
        # Section header
        story.append(Paragraph(section, styles['Heading2']))
        story.append(Spacer(1, 10))
        
        # Section content
        story.append(Paragraph(content, styles['Normal']))
        story.append(Spacer(1, 20))

    doc.build(story)
    buffer.seek(0)
    return buffer

def create_ppt(summary_data):
    """Create a PowerPoint presentation from the summary data."""
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "RFP Analysis Summary"
    subtitle.text = "Comprehensive Analysis Report"

    # Content slides
    for section, content in summary_data.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title.text = section
        content_placeholder.text = content

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def get_download_link(file_buffer, filename, file_type):
    """Generate a download link for the file."""
    b64 = base64.b64encode(file_buffer.getvalue()).decode()
    return f'<a href="data:application/{file_type};base64,{b64}" download="{filename}">Download {file_type.upper()}</a>'

def main():
    st.set_page_config(page_title="RFP Analyzer", layout="wide")
    st.title("RFP Document Analyzer")

    # File upload
    uploaded_file = st.file_uploader("Upload RFP Document (PDF)", type=['pdf'])
    
    if uploaded_file:
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
            tmp.write(uploaded_file.getvalue())
            tmp_path = tmp.name

        # Process document
        with st.spinner("Processing document..."):
            files = {"file": (uploaded_file.name, open(tmp_path, "rb"), "application/pdf")}
            response = requests.post(f"{API_URL}/analyze", files=files)
            if response.status_code == 200:
                st.success("Document processed successfully!")
            else:
                st.error("Error processing document")
                return

        # Query section
        st.header("Ask Questions")
        question = st.text_input("Enter your question about the RFP:")
        if question:
            with st.spinner("Generating answer..."):
                response = requests.post(f"{API_URL}/query", json={"question": question})
                if response.status_code == 200:
                    st.write("Answer:", response.json()["answer"])
                else:
                    st.error("Error getting answer")

        # Summary section
        st.header("Generate Summary")
        if st.button("Generate Summary"):
            with st.spinner("Generating summary..."):
                response = requests.get(f"{API_URL}/summary")
                if response.status_code == 200:
                    summary_data = response.json()
                    
                    # Display summary
                    st.subheader("Summary")
                    for section, content in summary_data.items():
                        with st.expander(section):
                            st.write(content)

                    # Export options
                    st.subheader("Export Summary")
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        if st.button("Export as PDF"):
                            pdf_buffer = create_pdf(summary_data)
                            st.markdown(get_download_link(pdf_buffer, "rfp_summary.pdf", "pdf"), unsafe_allow_html=True)
                    
                    with col2:
                        if st.button("Export as PowerPoint"):
                            ppt_buffer = create_ppt(summary_data)
                            st.markdown(get_download_link(ppt_buffer, "rfp_summary.pptx", "pptx"), unsafe_allow_html=True)

if __name__ == "__main__":
    main() 